from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from io import BytesIO
import streamlit as st


def create_ppt(dict_data: dict, paragraph: str, judul: str = 'Report pengolahan sentimen'):
    ppt = Presentation()

    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    # adding title
    title = slide.shapes.title
    title.text = judul

    # adding paragraph
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = paragraph
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    total = sum(dict_data.values())
    colors = [RGBColor(255, 0, 0), RGBColor(211, 211, 211), RGBColor(0, 128, 0)]
    # adding bar chart
    chart_data = CategoryChartData()
    chart_data.categories = list(dict_data.keys())
    chart_data.add_series('Sebaran Sentimen', (
        dict_data.get("Netral"),
        dict_data.get("Negatif"),
        dict_data.get("Positif")))
    x, y, cx, cy = Inches(0.5), Inches(3), Inches(3.75), Inches(3.75)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart

    for i, point in enumerate(chart.series[0].points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = colors[i]
        point.data_label.text = f'{dict_data[list(dict_data.keys())[i]]}'

    chart.has_label = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels

    data_labels.font.size = Pt(10)
    data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    # adding pie chart
    pie_chart_data = ChartData()
    pie_chart_data.categories = list(dict_data.keys())
    pie_chart_data.add_series('Sebaran', (
        dict_data.get("Netral") / total, dict_data.get("Negatif") / total, dict_data.get("Positif") / total))

    pie_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(5), Inches(3), Inches(4.5), Inches(4.5), pie_chart_data
    ).chart
    pie_chart.has_title = False
    pie_chart.has_legend = True
    pie_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    pie_chart.legend.include_in_layout = False

    pie_chart.plots[0].has_data_labels = True
    data_labels = pie_chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    # Mengatur warna pie chart
    for i, point in enumerate(pie_chart.plots[0].series[0].points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = colors[i]
    binary_output = BytesIO()
    ppt.save(binary_output)
    return binary_output.getvalue()
