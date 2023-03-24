import pptx
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.dml.color import RGBColor

# Create new PowerPoint presentation
prs = pptx.Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Create data for charts
sentiment_counts = [20, 30, 50]
labels = ['Positive', 'Neutral', 'Negative']
colors = [RGBColor(0, 176, 80), RGBColor(128, 128, 128), RGBColor(255, 0, 0)]

# Create pie chart data object
pie_data = ChartData()
pie_data.categories = labels
for i, count in enumerate(sentiment_counts):
    series = pie_data.add_series(labels[i], (count,))
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = colors[i]

# Create bar chart data object
bar_data = ChartData()
bar_data.categories = labels
bar_data.add_series('Sentiment Counts', sentiment_counts)

# Add pie chart to slide
x, y, cx, cy = Inches(1), Inches(2), Inches(4.5), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, x, y, cx, cy, pie_data
).chart

# Change chart title
chart.has_title = True
chart.chart_title.text_frame.text = "Sentiment Distribution"

# Add bar chart to slide
x, y, cx, cy = Inches(5.5), Inches(2), Inches(4.5), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, bar_data
).chart

# Change color of each column in the bar chart
for i, series in enumerate(chart.series):
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = colors[i]

# Change chart title
chart.has_title = True
chart.chart_title.text_frame.text = "Sentiment Counts"

# Save PowerPoint presentation
prs.save("sentiment_analysis.pptx")
